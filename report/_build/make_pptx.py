# -*- coding: utf-8 -*-
"""Genera la presentación (5 min) del Taller B4-T1 a partir del informe.
Tema académico-financiero coherente con el PDF. 16:9."""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

BASE = pathlib.Path(__file__).parent
ASSETS = BASE / "assets"
OUT = BASE.parent / "Taller_B4_T1_presentacion.pptx"

# --- paleta (idéntica al informe) ---
NAVY   = RGBColor(0x0e, 0x2a, 0x47)
NAVY2  = RGBColor(0x16, 0x3a, 0x61)
ACCENT = RGBColor(0x0f, 0x8b, 0x8d)
GOLD   = RGBColor(0xc8, 0xa2, 0x4a)
GOOD   = RGBColor(0x1b, 0x7a, 0x4b)
INK    = RGBColor(0x1c, 0x27, 0x33)
SLATE  = RGBColor(0x5b, 0x6b, 0x7b)
WHITE  = RGBColor(0xff, 0xff, 0xff)
SOFT   = RGBColor(0xf4, 0xf7, 0xfa)
LIGHT  = RGBColor(0xea, 0xf6, 0xf6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

HEAD = "Georgia"
BODY = "Calibri"


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def rrect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    """runs: list of (text, size, bold, color, font) or list of paragraphs (list of runs)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    paras = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space; p.space_after = Pt(4)
        for (t, sz, b, c, fn) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = fn
    return tb


def img_fit(s, path, x, y, max_w, max_h, center_x=True, center_y=True):
    """Inserta imagen escalada para caber en (max_w,max_h) conservando aspecto."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w; h = Emu(int(w / ar))
    if h > max_h:
        h = max_h; w = Emu(int(h * ar))
    px = x + (max_w - w) // 2 if center_x else x
    py = y + (max_h - h) // 2 if center_y else y
    s.shapes.add_picture(str(path), px, py, width=w, height=h)


def header(s, num, title, tag):
    rect(s, 0, 0, SW, Inches(1.0), NAVY)
    rect(s, 0, Inches(1.0), SW, Pt(3), GOLD)
    txt(s, Inches(0.5), Inches(0.12), Inches(0.7), Inches(0.7),
        [[(num, 30, True, GOLD, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.25), Inches(0.0), Inches(9.5), Inches(1.0),
        [[(title, 26, True, WHITE, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(10.6), Inches(0.0), Inches(2.4), Inches(1.0),
        [[(tag, 11, True, RGBColor(0x9f,0xd0,0xd6), BODY)]],
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def bullets(s, x, y, w, h, items, size=15, color=INK):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.08; p.space_after = Pt(9)
        r = p.add_run(); r.text = "▸ "; r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = ACCENT; r.font.name = BODY
        # it can be str or list of (text,bold)
        parts = [(it, False)] if isinstance(it, str) else it
        for (t, b) in parts:
            r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b
            r.font.color.rgb = color; r.font.name = BODY
    return tb


def kpi_row(s, y, items, w_each=Inches(2.9), gap=Inches(0.25), h=Inches(1.25), x0=None):
    n = len(items)
    total = w_each * n + gap * (n - 1)
    x = x0 if x0 is not None else (SW - total) // 2
    for (val, lab, col) in items:
        box = rrect(s, x, y, w_each, h, col)
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val; r.font.size = Pt(28); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = HEAD
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run(); r.text = lab; r.font.size = Pt(10.5); r.font.color.rgb = WHITE
        r.font.name = BODY
        x += w_each + gap


# ============ SLIDE 1 · PORTADA (académico-financiera) ============
s = slide()
# bandas navy arriba y abajo (fondo blanco por defecto)
rect(s, 0, 0, SW, Inches(0.95), NAVY)
rect(s, 0, Inches(6.78), SW, Inches(0.72), NAVY)
txt(s, Inches(0.6), 0, Inches(8.5), Inches(0.95),
    [[("Máster MIAX · IA aplicada a los mercados", 12.5, True, WHITE, BODY)]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(8.8), 0, Inches(3.9), Inches(0.95),
    [[("Bloque 4 · Tarea 1", 12.5, True, GOLD, BODY)]],
    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
# cuerpo
txt(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.4),
    [[("DOCUMENTO DE ENTREGA · JUNIO DE 2026", 11, False, SLATE, BODY)]])
txt(s, Inches(0.7), Inches(1.82), Inches(11.9), Inches(1.5),
    [[("Redes neuronales confiables para scoring de crédito:", 27, True, NAVY, HEAD)],
     [("equidad de género e incertidumbre", 27, True, NAVY, HEAD)]], space=1.06)
rect(s, Inches(0.72), Inches(3.5), Inches(2.0), Pt(4), GOLD)
tb = s.shapes.add_textbox(Inches(0.7), Inches(3.72), Inches(11.9), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Diseño, entrenamiento y auditoría sobre el dataset Home Credit Default Risk"
r.font.size = Pt(15); r.font.italic = True; r.font.color.rgb = SLATE; r.font.name = HEAD
txt(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.4),
    [[("Dataset:  ", 12, True, NAVY, BODY), ("Home Credit Default Risk — 307.511 × 122 variables", 12, False, INK, BODY)],
     [("Componentes:  ", 12, True, NAVY, BODY), ("capa custom (DTI) · FAIR loss (corr²) · Keras Tuner · MC-Dropout", 12, False, INK, BODY)],
     [("Variable sensible:  ", 12, True, NAVY, BODY), ("CODE_GENDER, no usada como entrada del modelo", 12, False, INK, BODY)]],
    space=1.35)
txt(s, Inches(0.7), Inches(5.98), Inches(11.9), Inches(0.32),
    [[("AUTORES", 9.5, True, SLATE, BODY)]])
txt(s, Inches(0.7), Inches(6.26), Inches(11.9), Inches(0.45),
    [[("Gonzalo de Ramón · Alonso Díaz · Oscar Romero", 16, True, NAVY, HEAD)]])
txt(s, Inches(0.6), Inches(6.78), Inches(12.1), Inches(0.72),
    [[("Taller B4-T1 — Diseño de Redes Confiables (Justicia e Incertidumbre)", 11, False,
       RGBColor(0xcf,0xe0,0xee), BODY)]], anchor=MSO_ANCHOR.MIDDLE)

# ============ SLIDE 2 · EL PROBLEMA ============
s = slide()
header(s, "▸", "El problema", "CONTEXTO")
bullets(s, Inches(0.6), Inches(1.35), Inches(6.4), Inches(5),
        [[("Scoring de perfiles con poco historial", False)],
         [("Desbalance 8,07 % (11,4:1) → métrica ", False), ("AUC-ROC", True)],
         [("Sesgo presente aunque el género no es input", False)],
         [("Proxy: ", False), ("EXT_SOURCE_1", True), (" (corr. −0,31 con género)", False)],
         [("Falta de fuentes externas predice impago → modelar incertidumbre", False)]], size=16)
# tres objetivos en cajas
y = Inches(1.5)
for (t, c) in [("PRECISIÓN", ACCENT), ("EQUIDAD", GOLD), ("INCERTIDUMBRE", GOOD)]:
    box = rrect(s, Inches(7.4), y, Inches(5.3), Inches(1.35), c)
    tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t; r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = HEAD
    y += Inches(1.55)
txt(s, Inches(7.4), Inches(6.25), Inches(5.3), Inches(0.8),
    [[("Los tres a la vez, en un mismo modelo.", 13, False, SLATE, BODY)]],
    align=PP_ALIGN.CENTER)

# ============ SLIDE 3 · TAREA 1 ============
s = slide()
header(s, "1", "Capa custom: restricción matemática", "CAPA CUSTOM")
bullets(s, Inches(0.6), Inches(1.35), Inches(6.2), Inches(4.2),
        [[("Capa Keras a medida ", False), ("DebtRatioSaturatingLayer", True)],
         [("Ratio: ", False), ("DTI = AMT_ANNUITY/(|AMT_INCOME|+ε)", True)],
         [("Saturación signed-power ", False), ("sign(x)·|x|ᵖ", True)],
         [("p entrenable ∈ [0,1;3], init 1 (identidad)", False)],
         [("Concatena → 14 features (primera capa)", False)],
         [("Aprende ", False), ("p ≈ 0,87", True)]], size=15.5)
box = rrect(s, Inches(0.6), Inches(5.7), Inches(6.2), Inches(1.45), RGBColor(0xfb,0xf6,0xea))
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Pt(10); tf.margin_right = Pt(10)
p = tf.paragraphs[0]
r = p.add_run()
r.text = ("No mejora AUC ni gap: opera sobre features ya transformadas. "
          "Aporte metodológico.")
r.font.size = Pt(13); r.font.color.rgb = INK; r.font.name = BODY
img_fit(s, ASSETS / "report__arquitectura_custom.png", Inches(7.1), Inches(1.3),
        Inches(5.9), Inches(5.9), center_x=True, center_y=True)

# ============ SLIDE 4 · TAREA 2 ============
s = slide()
header(s, "2", "FAIR loss: penalizar la dependencia", "FAIR LOSS")
# fórmula destacada
box = rrect(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.85), NAVY2)
tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "L = BCE(ŷ, TARGET) + λ · D(ŷ, CODE_GENDER)"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Consolas"
bullets(s, Inches(0.6), Inches(2.45), Inches(12.1), Inches(2.2),
        [[("D sobre la probabilidad; género solo en train, ", False), ("nunca input", True)],
         [("Tres medidas: corr² (lineal, O(n)), HSIC (kernel), MMD² (dos muestras)", False)],
         [("corr²", True), (": dependencia = desplazamiento de medias (binario) y O(n)", False)],
         [("Kernels exigen matrices n×n, inviables en validación", False)]], size=15.5)
box = rrect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.0), LIGHT)
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Pt(12)
p = tf.paragraphs[0]
r = p.add_run()
r.text = "No se borra el género (hay proxies): se penaliza su dependencia con la predicción."
r.font.size = Pt(15); r.font.color.rgb = INK; r.font.name = BODY
img_fit(s, ASSETS / "05_fair__score_medio_por_genero.png", Inches(0.6), Inches(6.05),
        Inches(12.1), Inches(1.25), center_x=True, center_y=False)

# ============ SLIDE 5 · TAREA 3 · PARETO ============
s = slide()
header(s, "3", "Frontera de Pareto: coste de la fairness", "KERAS TUNER")
img_fit(s, ASSETS / "06_tuner__pareto_3medidas.png", Inches(0.5), Inches(1.3),
        Inches(7.4), Inches(5.0), center_x=True, center_y=True)
kpi_items = [("−72 %", "group gap (5,615→1,568 pp)", GOLD),
             ("−0,59 pp", "AUC (vs base tuner)", ACCENT)]
# KPIs verticales a la derecha
y = Inches(1.5)
for (val, lab, col) in [("−72 %", "GROUP GAP  5,615 → 1,568 pp", GOLD),
                         ("−0,59 pp", "AUC vs base comparable", ACCENT),
                         ("−0,92 pp", "AUC vs base 03", NAVY2)]:
    box = rrect(s, Inches(8.1), y, Inches(4.7), Inches(1.15), col)
    tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = val; r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = HEAD
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = lab; r.font.size = Pt(10.5); r.font.color.rgb = WHITE; r.font.name = BODY
    y += Inches(1.3)
box = rrect(s, Inches(8.1), Inches(5.45), Inches(4.7), Inches(1.5), GOOD)
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Pt(10); tf.margin_right = Pt(10)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "−72 % de sesgo por menos de 1 pp de AUC."
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = BODY

# ============ SLIDE 6 · TAREA 4 · INCERTIDUMBRE ============
s = slide()
header(s, "4", "Incertidumbre: MC-Dropout (T=100)", "MC-DROPOUT")
img_fit(s, ASSETS / "07_incert__varianza_vs_n_ext_missing.png", Inches(0.5), Inches(1.3),
        Inches(6.4), Inches(5.0), center_x=True, center_y=True)
bullets(s, Inches(7.1), Inches(1.4), Inches(5.8), Inches(4),
        [[("Da clase + varianza; AUC 0,7346 (sin coste)", False)],
         [("Relación duda–EXT_SOURCE: global débil (ρ=+0,02)", False)],
         [("+47 % de varianza", True), (" sin fuentes externas (n=27)", False)],
         [("Buen vs mal pagador: varianza similar → duda = calidad de entrada", False)],
         [("Top 10 % incierto = más error → revisión humana", False)]], size=14.5)

# ============ SLIDE 7 · RESULTADOS (TABLA) ============
s = slide()
header(s, "5", "Resultados en test: base vs mejor FAIR", "EVALUACIÓN")
rows = [
    ("Modelo", "λ", "AUC test", "group gap", "ΔTPR", "ΔFPR"),
    ("base 03 (sin FAIR)", "—", "0,7437", "5,358", "—", "—"),
    ("base tuner (λ=0)", "0", "0,7404", "5,615", "8,282", "11,076"),
    ("★ mejor FAIR (corr²)", "5", "0,7345", "1,568", "1,817", "3,349"),
]
tbl_x, tbl_y = Inches(0.7), Inches(1.7)
tbl_w, tbl_h = Inches(11.9), Inches(2.9)
gtbl = s.shapes.add_table(len(rows), 6, tbl_x, tbl_y, tbl_w, tbl_h).table
widths = [Inches(3.9), Inches(1.0), Inches(1.9), Inches(2.1), Inches(1.5), Inches(1.5)]
for i, wd in enumerate(widths):
    gtbl.columns[i].width = wd
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.margin_left = Pt(6); cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val; r.font.name = BODY
        if ri == 0:
            r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        elif ri == 3:
            r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
        else:
            r.font.size = Pt(14); r.font.color.rgb = INK
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else SOFT
kpi_row(s, Inches(5.0),
        [("−72 %", "GROUP GAP", GOLD), ("−0,59 pp", "AUC (base comp.)", ACCENT),
         ("8,28→1,82", "ΔTPR (pp)", NAVY2), ("11,08→3,35", "ΔFPR (pp)", NAVY2)],
        w_each=Inches(2.9))
txt(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.7),
    [[("−72 % de sesgo y mejor equalized-odds por menos de 1 pp de AUC.", 14, False, SLATE, BODY)]],
    align=PP_ALIGN.CENTER)

# ============ SLIDE 8 · CONCLUSIONES ============
s = slide()
header(s, "6", "Conclusiones y defensa", "CIERRE")
data = [
    ("1 · Restricción y métrica",
     "Capa DTI (p entrenable) + penalización corr² sobre la probabilidad."),
    ("2 · Coste de la fairness",
     "−72 % de group gap por −0,59 pp de AUC (−0,92 vs base 03)."),
    ("3 · Incertidumbre y EXT_SOURCE",
     "Global débil, pero +47 % de varianza sin fuentes externas (n=27)."),
]
y = Inches(1.45)
for (h, b) in data:
    rect(s, Inches(0.6), y, Pt(5), Inches(1.25), ACCENT)
    txt(s, Inches(0.85), y, Inches(11.8), Inches(1.25),
        [[(h, 18, True, NAVY, HEAD)], [(b, 14.5, False, INK, BODY)]],
        anchor=MSO_ANCHOR.MIDDLE, space=1.05)
    y += Inches(1.42)
box = rrect(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.05), NAVY)
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Scoring + fairness + incertidumbre en un solo modelo."
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = HEAD

prs.save(str(OUT))
print("PPTX guardado:", OUT, "·", len(prs.slides.__iter__.__self__._sldIdLst), "diapositivas")
